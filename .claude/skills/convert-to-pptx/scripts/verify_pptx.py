#!/usr/bin/env python3
"""Structural sanity check on a produced .pptx.

Asserts the four properties this skill promises:
    1. No SVG parts inside the .pptx zip.
    2. No single picture covers the whole slide (the "screenshot trap").
    3. Each slide contains at least one editable text frame.
    4. Each slide contains multiple individually-selectable shapes.

Exits 0 on success. On failure, prints each problem and exits 1.

Usage:
    python verify_pptx.py output.pptx
"""
from __future__ import annotations

import argparse
import sys
import zipfile
from pathlib import Path

from pptx import Presentation

# A picture is considered "full-slide" if it covers >=95% of slide width AND
# >=95% of slide height. Some legitimate background photos cover the slide,
# but combined with other content; an offending deck typically has only the
# picture and maybe one or two captions, so we also flag slides with a
# full-slide picture and fewer than 3 total shapes.
FULL_SLIDE_RATIO = 0.95
MIN_SHAPES_WITH_FULL_PICTURE = 3


def find_svg_parts(pptx_path: Path) -> list[str]:
    with zipfile.ZipFile(pptx_path) as z:
        return [n for n in z.namelist() if n.lower().endswith(".svg")]


def check_slide(slide, slide_w: int, slide_h: int, idx: int) -> list[str]:
    problems: list[str] = []
    shapes = list(slide.shapes)

    if len(shapes) < 2:
        problems.append(
            f"slide {idx}: only {len(shapes)} shape(s) — expected multiple "
            f"individually-selectable shapes"
        )

    has_text_frame = any(
        getattr(s, "has_text_frame", False) and s.text_frame.text.strip()
        for s in shapes
    )
    if not has_text_frame:
        problems.append(
            f"slide {idx}: no text frame with content — text must be editable, "
            f"not flattened to an image"
        )

    for s in shapes:
        if s.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                w_ratio = (s.width or 0) / slide_w
                h_ratio = (s.height or 0) / slide_h
            except ZeroDivisionError:
                continue
            if (
                w_ratio >= FULL_SLIDE_RATIO
                and h_ratio >= FULL_SLIDE_RATIO
                and len(shapes) < MIN_SHAPES_WITH_FULL_PICTURE
            ):
                problems.append(
                    f"slide {idx}: a single picture covers the whole slide "
                    f"({w_ratio:.0%} x {h_ratio:.0%}) with only "
                    f"{len(shapes)} shape(s) — looks like a pasted screenshot"
                )

    return problems


def main() -> int:
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("pptx", type=Path)
    args = p.parse_args()

    if not args.pptx.exists():
        print(f"error: pptx not found: {args.pptx}", file=sys.stderr)
        return 2

    problems: list[str] = []

    svg_parts = find_svg_parts(args.pptx)
    if svg_parts:
        problems.append(
            "SVG parts found inside the .pptx (this skill forbids SVG): "
            + ", ".join(svg_parts)
        )

    prs = Presentation(args.pptx)
    slide_w = prs.slide_width or 0
    slide_h = prs.slide_height or 0
    if not slide_w or not slide_h:
        problems.append("slide width/height not set on the presentation")

    for i, slide in enumerate(prs.slides, start=1):
        problems.extend(check_slide(slide, slide_w, slide_h, i))

    if problems:
        print(f"verify_pptx: {len(problems)} problem(s) found in {args.pptx}")
        for prob in problems:
            print(f"  - {prob}")
        return 1

    print(f"verify_pptx: OK ({len(prs.slides)} slide(s))")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
